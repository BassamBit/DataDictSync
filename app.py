import io
import re
import numpy as np
import pandas as pd
import streamlit as st
from pptx import Presentation
from pptx.util import Pt
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

st.set_page_config(page_title="PowerPoint → Excel", page_icon="📥", layout="centered")
st.title("PowerPoint → Excel")
st.caption("ارفع ملف البوربوينت المعجم وسيتم استخراج البيانات إلى إكسل.")

# ========= Sidebar =========
with st.sidebar:
    st.header("الإعدادات")
    ROWS_PER_SLIDE = st.number_input("صفوف البيانات في كل شريحة", min_value=1, max_value=50, value=9)
    use_reversed = st.checkbox("الأعمدة معكوسة (reversed mapping)", value=True)
    OUTFILE_NAME = st.text_input("اسم ملف الإكسل", value="dictionary_output.xlsx")
    st.markdown("---")
    st.caption("الأعمدة المتوقعة في الجدول (9 أعمدة)")
    if use_reversed:
        st.markdown("""
| عمود الجدول | البيانات |
|---|---|
| 0 | AR Name+Def |
| 1 | Owner AR |
| 2 | Class AR |
| 3 | Code |
| 4 | Personal Data |
| 5 | Personal Data EN |
| 6 | Class EN |
| 7 | Owner EN |
| 8 | EN Name+Def |
        """)
    else:
        st.markdown("""
| عمود الجدول | البيانات |
|---|---|
| 0 | Code |
| 1 | EN Name+Def |
| 2 | Owner EN |
| 3 | Class EN |
| 4 | Personal Data |
| 5 | Personal Data EN |
| 6 | Class AR |
| 7 | Owner AR |
| 8 | AR Name+Def |
        """)

# ========= Upload =========
pptx_file = st.file_uploader("📎 ارفع ملف البوربوينت (.pptx)", type=["pptx"])

# ========= Helpers =========
def get_cell_text(cell):
    """Extract plain text from a table cell."""
    try:
        return cell.text_frame.text.strip()
    except:
        return ""

def split_name_def(text):
    """
    Split combined 'Name:\nDefinition' or 'Name\nDefinition' into two parts.
    Returns (name, definition).
    """
    if not text:
        return "", ""
    lines = text.splitlines()
    if len(lines) == 1:
        # No newline — could be name only or name: def on same line
        # Try splitting by first colon
        if ":" in lines[0]:
            parts = lines[0].split(":", 1)
            return parts[0].strip(), parts[1].strip()
        return lines[0].strip(), ""
    else:
        # First line = name (strip trailing colon)
        name = re.sub(r'[\s:：]+$', '', lines[0]).strip()
        definition = "\n".join(lines[1:]).strip()
        return name, definition

def extract_table_data(slide, num_cols=9):
    """Extract all data rows (skip header row 0) from first table in slide."""
    for sh in slide.shapes:
        if sh.has_table:
            table = sh.table
            rows = []
            for r in range(1, len(table.rows)):  # skip header
                row = []
                for c in range(min(num_cols, len(table.columns))):
                    row.append(get_cell_text(table.cell(r, c)))
                # Pad if fewer columns
                while len(row) < num_cols:
                    row.append("")
                rows.append(row)
            return rows
    return []

def is_empty_row(row):
    return all(not str(v).strip() for v in row)

# ========= Column mapping =========
# REVERSED mapping (default): col index in table → meaning
# 0: AR Name+Def, 1: Owner AR, 2: Class AR, 3: Code,
# 4: Personal Data, 5: Personal Data EN, 6: Class EN, 7: Owner EN, 8: EN Name+Def

REVERSED_MAP = {
    3: "code",
    8: "en_name_def",
    7: "owner_en",
    6: "class_en",
    4: "personal_data",
    5: "personal_data_en",
    2: "class_ar",
    1: "owner_ar",
    0: "ar_name_def",
}

NATURAL_MAP = {
    0: "code",
    1: "en_name_def",
    2: "owner_en",
    3: "class_en",
    4: "personal_data",
    5: "personal_data_en",
    6: "class_ar",
    7: "owner_ar",
    8: "ar_name_def",
}

# ========= Run =========
run = st.button("استخراج البيانات ✅", type="primary", use_container_width=True)

if run:
    if pptx_file is None:
        st.error("الرجاء رفع ملف البوربوينت أولاً.")
        st.stop()

    prs = Presentation(io.BytesIO(pptx_file.read()))
    st.info(f"عدد الشرائح: {len(prs.slides)}")

    col_map = REVERSED_MAP if use_reversed else NATURAL_MAP
    all_records = []

    for slide_idx, slide in enumerate(prs.slides):
        rows = extract_table_data(slide, num_cols=9)
        for row in rows:
            if is_empty_row(row):
                continue

            record = {}
            for col_idx, meaning in col_map.items():
                if col_idx < len(row):
                    record[meaning] = row[col_idx]
                else:
                    record[meaning] = ""

            # Split combined name+def fields
            en_name, en_def = split_name_def(record.get("en_name_def", ""))
            ar_name, ar_def = split_name_def(record.get("ar_name_def", ""))

            all_records.append({
                "Code":                record.get("code", ""),
                "English Name":        en_name,
                "English Definition":  en_def,
                "Classification EN":   record.get("class_en", ""),
                "Personal Data":       record.get("personal_data", ""),
                "Arabic Name":         ar_name,
                "Arabic Definition":   ar_def,
                "Classification AR":   record.get("class_ar", ""),
                "Data Owner EN":       record.get("owner_en", ""),
                "Data Owner AR":       record.get("owner_ar", ""),
            })

    if not all_records:
        st.error("لم يتم العثور على بيانات في الملف. تأكد من الإعدادات.")
        st.stop()

    df_out = pd.DataFrame(all_records)
    # Remove completely empty rows
    df_out = df_out[~df_out.apply(lambda r: all(str(v).strip() == "" for v in r), axis=1)]
    df_out = df_out.reset_index(drop=True)

    st.success(f"✅ تم استخراج {len(df_out)} سجل من {len(prs.slides)} شريحة.")
    st.dataframe(df_out, use_container_width=True)

    # ========= Build Excel =========
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Dictionary"

    # Header style
    header_fill = PatternFill("solid", start_color="1F4E79")
    header_font = Font(bold=True, color="FFFFFF", size=11, name="Arial")
    center = Alignment(horizontal="center", vertical="center", wrap_text=True)
    left_align = Alignment(horizontal="left", vertical="center", wrap_text=True)
    right_align = Alignment(horizontal="right", vertical="center", wrap_text=True)
    thin = Side(style="thin", color="CCCCCC")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    headers = list(df_out.columns)
    for c_idx, h in enumerate(headers, 1):
        cell = ws.cell(row=1, column=c_idx, value=h)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = center
        cell.border = border

    # Data rows
    alt_fill = PatternFill("solid", start_color="EBF2FA")
    for r_idx, row in df_out.iterrows():
        fill = alt_fill if r_idx % 2 == 1 else PatternFill("solid", start_color="FFFFFF")
        for c_idx, val in enumerate(row, 1):
            cell = ws.cell(row=r_idx+2, column=c_idx, value=str(val) if val else "")
            cell.font = Font(size=11, name="Arial")
            cell.border = border
            cell.fill = fill
            # Arabic columns: right align
            col_name = headers[c_idx-1]
            if "Arabic" in col_name or "AR" in col_name:
                cell.alignment = right_align
            elif col_name == "Code":
                cell.alignment = center
            else:
                cell.alignment = left_align

    # Column widths
    col_widths = {
        "Code": 12,
        "English Name": 25,
        "English Definition": 45,
        "Classification EN": 18,
        "Personal Data": 16,
        "Arabic Name": 25,
        "Arabic Definition": 45,
        "Classification AR": 18,
        "Data Owner EN": 22,
        "Data Owner AR": 22,
    }
    for c_idx, h in enumerate(headers, 1):
        ws.column_dimensions[get_column_letter(c_idx)].width = col_widths.get(h, 20)

    ws.row_dimensions[1].height = 30
    ws.freeze_panes = "A2"

    # Save
    out_buf = io.BytesIO()
    wb.save(out_buf)
    out_buf.seek(0)

    st.download_button(
        "⬇️ تحميل ملف الإكسل",
        data=out_buf,
        file_name=OUTFILE_NAME or "dictionary_output.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        use_container_width=True
    )

st.markdown("---")
st.caption("يستخرج البيانات من جداول البوربوينت ويفصل الاسم عن التعريف تلقائياً.")
